"""文档读取器 - 支持多种格式"""

from __future__ import annotations

from pathlib import Path
from typing import Optional


class DocumentReader:
    """文档读取器，支持多种文件格式"""

    SUPPORTED_EXTENSIONS = {
        ".txt": "read_text",
        ".md": "read_markdown",
        ".pdf": "read_pdf",
        ".docx": "read_docx",
        ".pptx": "read_pptx",
        ".xlsx": "read_xlsx",
        ".html": "read_html",
        ".htm": "read_html",
        ".csv": "read_csv",
    }

    @staticmethod
    def read(file_path: str) -> tuple[str, dict]:
        """读取文档内容，返回 (文本内容, 元数据)"""
        ext = Path(file_path).suffix.lower()
        if ext not in DocumentReader.SUPPORTED_EXTENSIONS:
            raise ValueError(f"不支持的文件格式: {ext}，支持: {list(DocumentReader.SUPPORTED_EXTENSIONS.keys())}")

        reader_method = DocumentReader.SUPPORTED_EXTENSIONS[ext]
        return getattr(DocumentReader, reader_method)(file_path)

    @staticmethod
    def read_text(file_path: str) -> tuple[str, dict]:
        """读取纯文本文件"""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return text, {"source": file_path, "format": "txt"}

    @staticmethod
    def read_markdown(file_path: str) -> tuple[str, dict]:
        """读取 Markdown 文件"""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return text, {"source": file_path, "format": "markdown"}

    @staticmethod
    def read_pdf(file_path: str) -> tuple[str, dict]:
        """读取 PDF 文件，使用 PyMuPDF 提取文本"""
        try:
            import pypdf
        except ImportError:
            raise ImportError("请安装 pypdf: pip install pypdf")

        text_parts = []
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"--- 第 {page_num + 1} 页 ---\n{page_text}")

        full_text = "\n\n".join(text_parts)
        return full_text, {
            "source": file_path,
            "format": "pdf",
            "page_count": len(reader.pages),
        }

    @staticmethod
    def read_docx(file_path: str) -> tuple[str, dict]:
        """读取 Word 文档"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")

        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                text_parts.append(row_text)

        full_text = "\n\n".join(text_parts)
        return full_text, {"source": file_path, "format": "docx"}

    @staticmethod
    def read_pptx(file_path: str) -> tuple[str, dict]:
        """读取 PowerPoint 文件"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- 幻灯片 {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            text_parts.append("\n".join(slide_texts))

        full_text = "\n\n".join(text_parts)
        return full_text, {"source": file_path, "format": "pptx", "slide_count": len(prs.slides)}

    @staticmethod
    def read_xlsx(file_path: str) -> tuple[str, dict]:
        """读取 Excel 文件"""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("请安装 openpyxl: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path, read_only=True)
        text_parts = []
        sheet_names = wb.sheetnames
        for sheet_name in sheet_names:
            ws = wb[sheet_name]
            sheet_texts = [f"--- 工作表: {sheet_name} ---"]
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(cell) if cell is not None else "" for cell in row
                )
                if row_text.strip():
                    sheet_texts.append(row_text)
            text_parts.append("\n".join(sheet_texts))

        wb.close()
        full_text = "\n\n".join(text_parts)
        return full_text, {
            "source": file_path,
            "format": "xlsx",
            "sheet_count": len(sheet_names),
        }

    @staticmethod
    def read_html(file_path: str) -> tuple[str, dict]:
        """读取 HTML 文件"""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            html_content = f.read()

        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html_content, "html.parser")
            # 移除 script 和 style 标签
            for tag in soup(["script", "style"]):
                tag.decompose()
            text = soup.get_text(separator="\n")
            # 清理多余空行
            lines = [line.strip() for line in text.split("\n") if line.strip()]
            full_text = "\n".join(lines)
        except ImportError:
            # 没有BeautifulSoup时做简单处理
            import re
            full_text = re.sub(r"<[^>]+>", " ", html_content)
            full_text = re.sub(r"\s+", " ", full_text).strip()

        return full_text, {"source": file_path, "format": "html"}

    @staticmethod
    def read_csv(file_path: str) -> tuple[str, dict]:
        """读取 CSV 文件"""
        import csv
        text_parts = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                text_parts.append(" | ".join(row))
        full_text = "\n".join(text_parts)
        return full_text, {"source": file_path, "format": "csv"}
